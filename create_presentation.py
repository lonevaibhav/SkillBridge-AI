"""
Generate PowerPoint presentation for SkillBridge AI project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create a professional PowerPoint presentation for SkillBridge AI"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    PRIMARY_COLOR = RGBColor(0, 102, 204)  # Blue
    ACCENT_COLOR = RGBColor(255, 102, 0)   # Orange
    TEXT_COLOR = RGBColor(51, 51, 51)      # Dark gray
    
    def add_title_slide(title, subtitle):
        """Add a title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = ACCENT_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(title, content_points):
        """Add a content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # Add horizontal line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
        line.line.color.rgb = ACCENT_COLOR
        line.line.width = Pt(3)
        
        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.2))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(content_points):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)
    
    # Slide 1: Title Slide
    add_title_slide("SkillBridge AI", "AI-Powered Career & Placement Assistant")
    
    # Slide 2: Problem Statement
    add_content_slide(
        "Problem Statement",
        [
            "• Students struggle to identify skill gaps for target job roles",
            "• Resume screening is time-consuming and subjective",
            "• Limited personalized career guidance and recommendations",
            "• Difficulty matching student profiles with relevant opportunities",
            "• Need for data-driven insights on career development"
        ]
    )
    
    # Slide 3: Solution Overview
    add_content_slide(
        "Solution: SkillBridge AI",
        [
            "• AI-powered platform for career guidance and placement",
            "• Automated resume parsing and skill extraction",
            "• Intelligent job-to-profile matching engine",
            "• Personalized skill gap analysis and recommendations",
            "• RAG-based approach with local LLM (Ollama)",
            "• Scalable architecture with vector search capabilities"
        ]
    )
    
    # Slide 4: Tech Stack
    add_content_slide(
        "Technology Stack",
        [
            "Frontend: Next.js 14 + TypeScript + Vite",
            "Backend: FastAPI + Python",
            "AI/ML: Ollama (Local LLM) + FAISS (Vector Search)",
            "Database: MongoDB + Motor (Async Driver)",
            "ML Models: Resume Parser, Skill Matcher, Experience Extractor",
            "Deployment: Docker & Docker Compose"
        ]
    )
    
    # Slide 5: Architecture
    add_content_slide(
        "System Architecture",
        [
            "Frontend Layer: Modern React UI for user interactions",
            "API Layer: FastAPI REST endpoints for all operations",
            "ML Layer: Python-based ML models for intelligence",
            "LLM Integration: Ollama for advanced text generation",
            "Vector DB: FAISS for semantic skill matching",
            "Data Layer: MongoDB for persistent storage"
        ]
    )
    
    # Slide 6: Key Features
    add_content_slide(
        "Key Features",
        [
            "✓ Resume Upload & Parsing: Extract skills, experience, contact info",
            "✓ Skill Analysis: AI-powered insights on skills and gaps",
            "✓ Job Matching: Find opportunities based on profile",
            "✓ Match Analysis: Detailed compatibility scoring",
            "✓ Recommendations: Personalized career development advice",
            "✓ Experience Tracking: Automated career progression analysis"
        ]
    )
    
    # Slide 7: ML Models
    add_content_slide(
        "Machine Learning Models",
        [
            "Resume Parser: Extracts text from PDF/DOCX/TXT files",
            "Skill Matcher: Semantic matching using embeddings",
            "Experience Extractor: Analyzes work history & progression",
            "Skill Embeddings: sentence-transformers for semantic search",
            "Ollama/Mistral: Local LLM for advanced NLP tasks",
            "FAISS: Vector database for efficient similarity search"
        ]
    )
    
    # Slide 8: Data Flow
    add_content_slide(
        "Data Processing Pipeline",
        [
            "1. User uploads resume (PDF/DOCX)",
            "2. Backend parses and extracts structured data",
            "3. ML models process: skills, experience, level",
            "4. Vector embeddings created for semantic search",
            "5. Data stored in MongoDB for persistence",
            "6. Real-time matching against job database using FAISS"
        ]
    )
    
    # Slide 9: Use Cases
    add_content_slide(
        "Use Cases",
        [
            "Students: Identify career paths and required skills",
            "Job Seekers: Match profiles with opportunities",
            "Recruiters: Quickly screen candidates by skills",
            "Career Counselors: Data-driven placement guidance",
            "Institutions: Analyze student placement readiness",
            "HR Teams: Automate initial screening process"
        ]
    )
    
    # Slide 10: Demo Results
    add_content_slide(
        "Expected Outcomes",
        [
            "✓ Accurate skill extraction from resumes (95%+ accuracy)",
            "✓ Intelligent job matching with relevance scoring",
            "✓ Personalized recommendations based on profile",
            "✓ Reduced manual resume screening time by 70%+",
            "✓ Improved placement rate through better matching",
            "✓ Scalable to handle 1000+ concurrent users"
        ]
    )
    
    # Slide 11: Future Enhancements
    add_content_slide(
        "Future Enhancements",
        [
            "• Real-time interview preparation with AI coaching",
            "• LinkedIn profile integration for auto-parsing",
            "• Advanced analytics dashboard for trends",
            "• Mobile app for on-the-go career guidance",
            "• Integration with job portals (LinkedIn, Indeed, etc.)",
            "• Multi-language support for global reach"
        ]
    )
    
    # Slide 12: Team & Implementation
    add_content_slide(
        "Implementation Highlights",
        [
            "• Full-stack development with modern frameworks",
            "• Custom ML models for domain-specific extraction",
            "• Async processing for high performance",
            "• Docker containerization for easy deployment",
            "• RAG architecture for intelligent responses",
            "• Comprehensive error handling and logging"
        ]
    )
    
    # Slide 13: Conclusion
    add_content_slide(
        "Conclusion",
        [
            "• SkillBridge AI revolutionizes career guidance through AI",
            "• Bridges gap between education and employment",
            "• Leverages cutting-edge ML and LLM technologies",
            "• Scalable, production-ready solution",
            "• Addresses real-world placement challenges",
            "• Opens opportunities for continuous improvement"
        ]
    )
    
    # Slide 14: Thank You
    add_title_slide("Thank You!", "Questions & Discussion")
    
    # Save presentation
    output_path = "SkillBridge_AI_Presentation.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
